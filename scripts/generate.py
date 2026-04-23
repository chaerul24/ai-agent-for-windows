from pptx import Presentation
from pptx.util import Pt
import json
import sys
import re
import requests
import os
import time

# ===== CONFIG =====
UNSPLASH_KEY = "DeDraJwRBF10LUXAUl3iBrclyVcIK6mmYLwSVZzBzkE"

MAX_POINTS = 5
MAX_LENGTH = 80

# ===== SAFE FILENAME =====
def safe_filename(name):
    name = name.lower()
    name = re.sub(r'[^a-z0-9]', '_', name)
    name = re.sub(r'_+', '_', name)
    return name.strip('_')

# ===== CLEAN TEXT =====
def clean_text(text, max_len):
    if not text:
        return ""

    text = str(text).strip()

    if len(text) > max_len:
        return text[:max_len - 3] + "..."

    return text

# ===== CLEAN POINTS =====
def clean_points(points):
    if not points:
        return ["Konten tidak tersedia"]

    points = points[:MAX_POINTS]

    clean = []
    for p in points:
        p = clean_text(p, MAX_LENGTH)
        clean.append(p)

    return clean

# ===== DOWNLOAD IMAGE FROM UNSPLASH =====
def fetch_image(query):
    try:
        if not UNSPLASH_KEY:
            return None

        url = f"https://api.unsplash.com/photos/random?query={query}&orientation=landscape"

        res = requests.get(url, headers={
            "Authorization": f"Client-ID {UNSPLASH_KEY}"
        })

        data = res.json()

        if "urls" not in data:
            return None

        img_url = data["urls"]["regular"]
        img_data = requests.get(img_url).content

        path = f"temp/{int(time.time()*1000)}.jpg"

        with open(path, "wb") as f:
            f.write(img_data)

        return path

    except Exception as e:
        print("⚠️ Image error:", str(e))
        return None

# ===== LOAD DATA =====
raw = sys.argv[1]
raw = raw.replace('\\"', '"')

data = json.loads(raw)

# ===== LOAD TEMPLATE =====
prs = Presentation("templates/template.pptx")

slides = prs.slides

# ===== PROCESS SLIDES =====
for i, slide_data in enumerate(data["slides"]):
    if i >= len(slides):
        break

    slide = slides[i]

    # ===== TITLE =====
    title_text = clean_text(slide_data.get("title", ""), 60)

    if slide.shapes.title:
        slide.shapes.title.text = title_text

    # ===== CONTENT =====
    points = clean_points(slide_data.get("points", []))

    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for idx, p in enumerate(points):
            if idx == 0:
                tf.text = p
            else:
                para = tf.add_paragraph()
                para.text = p
                para.level = 0

    # ===== IMAGE =====
    image_query = slide_data.get("image")

    if image_query:
        img_path = fetch_image(image_query)

        if img_path:
            try:
                slide.shapes.add_picture(
                    img_path,
                    left=prs.slide_width * 0.6,
                    top=prs.slide_height * 0.3,
                    width=prs.slide_width * 0.3
                )
            except Exception as e:
                print("⚠️ Gagal insert gambar:", str(e))

# ===== SAVE FILE =====
file_name = safe_filename(data.get("title", "presentation")) + ".pptx"

prs.save(file_name)

print(f"✅ Saved: {file_name}")